#!/usr/bin/env python3
"""
Build a .pptx corpus with machine-readable ground truth.

Every shape is named with the verdict it is *designed* to produce:

    FIT_*        the text should fit
    OVER_*       the text should clearly overflow
    OVERLAP_*    the shape should intersect another
    OFFCANVAS_*  the shape should extend past the slide edge

Designed, not asserted. What the shapes actually do is settled by rendering
them, which is what calibrate_pptx.py does. Naming the intent separately from
the measurement is the only way to tell "my measurement is wrong" apart from
"my fixture is wrong" - and two labels here were wrong on the first pass, which
is exactly why the separation is worth the trouble.
"""
from __future__ import annotations

import os

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Emu, Pt

OUT = os.environ.get("DI_PPTX", "../corpus/deck.pptx")

# 16:9 at the usual 13.333 x 7.5 inches
SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)

LOREM = ("The Supplier shall maintain professional indemnity insurance of not "
         "less than five million euro per claim and in the aggregate, and shall "
         "provide evidence of such cover on each anniversary of the Effective "
         "Date throughout the Term and for six years thereafter.")


def _box(slide, name, left, top, width, height, text, *,
         size=18, bold=False, font="Calibri", wrap=True,
         autofit=MSO_AUTO_SIZE.NONE, anchor=MSO_ANCHOR.TOP,
         insets=None):
    shape = slide.shapes.add_textbox(Emu(left), Emu(top), Emu(width), Emu(height))
    shape.name = name
    tf = shape.text_frame
    tf.word_wrap = wrap
    tf.auto_size = autofit
    tf.vertical_anchor = anchor
    if insets is not None:
        tf.margin_left, tf.margin_top, tf.margin_right, tf.margin_bottom = insets
    tf.text = text
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.name = font
    return shape


def build() -> str:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    # ---------------------------------------------------------------- slide 1
    # The core question: does the text fit the box? Same box, three lengths.
    s = prs.slides.add_slide(blank)
    _box(s, "FIT_one_short_line", 500000, 400000, 3500000, 900000,
         "Commercial terms", size=20, bold=True)
    _box(s, "FIT_wrapped_paragraph", 500000, 1500000, 3500000, 2200000,
         LOREM[:120], size=14)
    _box(s, "OVER_paragraph_in_small_box", 4500000, 400000, 3200000, 700000,
         LOREM, size=14)
    _box(s, "OVER_huge_type_tiny_box", 8200000, 400000, 3000000, 600000,
         "Quarterly revenue by region", size=40, bold=True)
    _box(s, "FIT_generous_box", 4500000, 1500000, 6700000, 2600000, LOREM, size=14)

    # ---------------------------------------------------------------- slide 2
    # Wrap off: a single line runs past the right edge of the box regardless
    # of how tall the box is. A different failure from vertical overflow.
    s = prs.slides.add_slide(blank)
    _box(s, "OVER_nowrap_single_line", 500000, 500000, 2500000, 1200000,
         "Effective Date means the date of last signature",
         size=16, wrap=False)
    _box(s, "FIT_nowrap_short", 500000, 2200000, 2500000, 600000,
         "MSA-2026-0417", size=16, wrap=False)
    # Non-default insets: the usable area is smaller than the shape, which is a
    # classic source of "but it looks like it fits" disagreements.
    _box(s, "OVER_fat_insets", 4000000, 500000, 3000000, 1000000,
         LOREM[:90], size=14,
         insets=(Emu(457200), Emu(457200), Emu(457200), Emu(457200)))
    _box(s, "FIT_zero_insets", 4000000, 2000000, 3000000, 1000000,
         LOREM[:90], size=14, insets=(0, 0, 0, 0))
    # Labelled TIGHT_ at first. Both the measurement and LibreOffice put it at
    # 2 lines in a box with room for 4, so the label was wrong, not the checker.
    _box(s, "FIT_two_lines_of_four", 7500000, 500000, 4000000, 1000000,
         "The Supplier shall deliver the services with reasonable skill",
         size=14)

    # ---------------------------------------------------------------- slide 3
    # Geometry problems that need no text measurement at all.
    s = prs.slides.add_slide(blank)
    _box(s, "OVERLAP_lower_left", 1000000, 1000000, 3000000, 1500000,
         "Discovery phase", size=16)
    _box(s, "OVERLAP_upper_right", 2500000, 1800000, 3000000, 1500000,
         "Delivery phase", size=16)
    _box(s, "FIT_clear_of_others", 7000000, 1000000, 3000000, 1200000,
         "Handover", size=16)
    _box(s, "OFFCANVAS_right", 11000000, 3000000, 3000000, 800000,
         "Appendix reference", size=14)
    _box(s, "OFFCANVAS_bottom", 1000000, 6500000, 3000000, 900000,
         "Confidential draft", size=14)

    # ---------------------------------------------------------------- slide 4
    # Autofit variants. PowerPoint shrinks text (normAutofit) or grows the box
    # (spAutoFit); both change what "overflow" means, and a checker that
    # ignores them will cry wolf.
    s = prs.slides.add_slide(blank)
    _box(s, "AUTOFIT_shrink_text", 500000, 500000, 3000000, 900000, LOREM[:150],
         size=18, autofit=MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE)
    _box(s, "AUTOFIT_grow_shape", 4000000, 500000, 3000000, 900000, LOREM[:150],
         size=18, autofit=MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT)
    _box(s, "OVER_no_autofit_same_text", 7500000, 500000, 3000000, 900000,
         LOREM[:150], size=18, autofit=MSO_AUTO_SIZE.NONE)
    _box(s, "FIT_middle_anchored", 500000, 2200000, 3000000, 1500000,
         "Centred vertically", size=16, anchor=MSO_ANCHOR.MIDDLE)

    # ---------------------------------------------------------------- slide 5
    # Mixed run sizes in one paragraph, and a hard line break. The tallest run
    # sets the line height, so measuring per-paragraph at one size is wrong.
    s = prs.slides.add_slide(blank)
    shape = s.shapes.add_textbox(Emu(500000), Emu(500000),
                                 Emu(4000000), Emu(1000000))
    # Also mislabelled at first. A 32pt run plus a 12pt run in a 71pt box does
    # fit - the two engines disagree only on where the first line breaks, which
    # is the 0.6%-margin case described in pptx_checks.BORDERLINE.
    shape.name = "FIT_mixed_run_sizes"
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    p = tf.paragraphs[0]
    for text, size, bold in (("EUR 44,500", 32, True),
                             (" per annum, payable quarterly in arrears", 12, False)):
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.name = "Calibri"

    _box(s, "OVER_hard_breaks", 5000000, 500000, 3000000, 900000,
         "Milestone one\nMilestone two\nMilestone three\nMilestone four\n"
         "Milestone five", size=16)
    _box(s, "FIT_hard_breaks_room", 8500000, 500000, 3000000, 2400000,
         "Milestone one\nMilestone two\nMilestone three\nMilestone four\n"
         "Milestone five", size=16)
    _box(s, "FIT_bold_narrow", 500000, 2000000, 3000000, 1000000,
         "Confidential", size=24, bold=True)
    # A face that is not installed anywhere here, so measurement must degrade
    # to an estimate and say so rather than pretend.
    _box(s, "FIT_unknown_font", 5000000, 2000000, 3500000, 1000000,
         "Segoe UI sample text", size=18, font="Segoe UI")

    os.makedirs(os.path.dirname(OUT) or ".", exist_ok=True)
    prs.save(OUT)
    return OUT


def ground_truth(path: str) -> dict[str, str]:
    """The verdict each shape's name declares. Intent, not measurement."""
    from pptx import Presentation as P
    out: dict[str, str] = {}
    for i, slide in enumerate(P(path).slides, 1):
        for shape in slide.shapes:
            out[f"slide{i}/{shape.name}"] = shape.name.split("_", 1)[0]
    return out


if __name__ == "__main__":
    p = build()
    gt = ground_truth(p)
    from collections import Counter
    print(f"built {p}  ({os.path.getsize(p)} bytes, {len(gt)} shapes)")
    for kind, n in sorted(Counter(gt.values()).items()):
        print(f"  {kind:10} {n}")
