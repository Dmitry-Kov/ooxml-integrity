#!/usr/bin/env python3
"""Draw a visible outline around every text box, for checking against a renderer.

Why this exists. The checker's central claim is "this text needs more height
than its box gives it". A screenshot of the deck cannot confirm or refute that,
because the box has no border: overflowing text and merely long text look
identical. With an outline on every shape the question becomes something a
person can answer by looking - does the text cross the line or not - with no
measuring and no trust in my arithmetic.

The outline is decoration (`a:ln` on the shape's fill properties). It changes no
inset, no wrap, no autofit, so the layout being judged is the same layout. That
is asserted rather than assumed: --check re-runs the checks on the copy and
compares them to the original.

    python outline_deck.py                 # ../corpus/deck.pptx -> deck_outlined.pptx
    python outline_deck.py --check         # and verify the findings are unchanged
"""
from __future__ import annotations

import argparse
import collections
import os
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt

IN = os.environ.get("DI_PPTX", "../corpus/deck.pptx")
OUT = os.environ.get("DI_PPTX_OUTLINED", "../corpus/deck_outlined.pptx")

#: Magenta, because nothing in the deck is magenta. Thin, so the line itself
#: cannot be mistaken for the text's own edge.
LINE_RGB = RGBColor(0xE0, 0x00, 0x90)
LINE_PT = 0.75


def outline(src: str = IN, dst: str = OUT) -> str:
    prs = Presentation(src)
    touched = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            line = shape.line
            line.color.rgb = LINE_RGB
            line.width = Pt(LINE_PT)
            touched += 1
    os.makedirs(os.path.dirname(dst) or ".", exist_ok=True)
    prs.save(dst)
    print(f"outlined {touched} text shapes -> {dst}")
    return dst


def same_findings(a: str, b: str) -> bool:
    from ooxml_integrity import check_pptx

    def codes(p):
        return collections.Counter(f"{f.code} {f.where}" for f in check_pptx(p))

    ca, cb = codes(a), codes(b)
    if ca == cb:
        print(f"both report the same {sum(ca.values())} findings - the outline "
              f"is decoration, not a layout change")
        return True
    print("the outline changed the verdicts, so the copy cannot stand in for "
          "the original:", file=sys.stderr)
    for key in sorted(set(ca) | set(cb)):
        if ca.get(key, 0) != cb.get(key, 0):
            print(f"  {key}: original {ca.get(key, 0)}, outlined {cb.get(key, 0)}",
                  file=sys.stderr)
    return False


if __name__ == "__main__":
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--check", action="store_true",
                    help="verify the copy reports exactly what the original does")
    args = ap.parse_args()
    dst = outline()
    if args.check and not same_findings(IN, dst):
        sys.exit(1)
