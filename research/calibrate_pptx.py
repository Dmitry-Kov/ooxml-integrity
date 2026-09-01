#!/usr/bin/env python3
"""
Calibrate the text measurement against a real renderer.

The question: when the checker says a paragraph wraps to 7 lines with a line
pitch of 17.1pt, how close is that to what a renderer actually does? Without
this number every overflow verdict is an assertion.

Method, and why it is shaped this way:

* **One shape per rendered deck.** An earlier version tried to attribute text
  lines to shapes by geometry on a full slide, and got it wrong wherever two
  shapes shared a horizontal band - producing nonsense like a 767pt "measured"
  height for a one-line label. Rendering each shape alone removes the
  attribution problem entirely.
* **Compare line count and line pitch, not total height.** pdfplumber reports a
  character's height from the font matrix, so a single rendered line comes back
  as the font size rather than the line box. Comparing that against a computed
  line height produced a constant +22% error that was an artefact of comparing
  two different quantities. Line count and baseline-to-baseline pitch are both
  unambiguous, and together they determine the height.
* **Autofit off, box made huge.** The point is to measure how text lays out, not
  how a renderer clips it. Each probe puts the shape's text in a box wide enough
  to preserve the original wrapping and tall enough that nothing is dropped.

The renderer is a parameter, not a constant. `--renderer` accepts `soffice`
(LibreOffice), `x2t` (the converter inside ONLYOFFICE Desktop Editors, found
automatically or given as a path) or any explicit binary path, and `--json`
writes the per-shape numbers out so several renderers can be put side by side
rather than compared from memory.

What agreement here does and does not show: a renderer agreeing with the model
means the model is sane for that renderer. PowerPoint was checked separately and
by eye, on the outlined deck - see docs/powerpoint-validation.md. Where two
renderers disagree with each other, at most one of them can match the model, and
the interesting question becomes which.
"""
from __future__ import annotations

import argparse
import glob
import json
import os
import shutil
import statistics
import subprocess
import sys
from collections import defaultdict

import pdfplumber
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Emu, Pt

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "src"))

from ooxml_integrity.fonts import EMU_PER_POINT  # noqa: E402
from ooxml_integrity.pptx_layout import (  # noqa: E402
    DRAWINGML_LINE_SPACING, layout_shape, read_deck,
)

DECK = os.environ.get("DI_PPTX", "../corpus/deck.pptx")
WORK = "/tmp/pptx-calib"


def probe_deck(shape, out_path: str) -> None:
    """One slide, one text box: same text, same width, same font, room to spare."""
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    _add_probe_slide(prs, shape)
    prs.save(out_path)


def probe_book(shapes, out_path: str, manifest_path: str) -> list[str]:
    """All the probes in one deck, one shape per slide, plus a page manifest.

    This exists for renderers with no usable command line. ONLYOFFICE ships the
    same `x2t` its server products use, but driven from outside the app bundle
    DoctRenderer fails to open the document (`<error code="open"/>` and a JS
    TypeError), because the sdkjs resources and font cache it expects are not
    set up. Rather than reverse-engineer that, this writes one file a person can
    export to PDF from the GUI in a single action - and because each probe gets
    its own slide, PDF page N is shape N with no attribution guesswork. That was
    the bug that produced a 767pt "measured" height for a one-line label when an
    earlier version tried to attribute lines to shapes on a shared slide.
    """
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    order: list[str] = []
    for shape in shapes:
        _add_probe_slide(prs, shape)
        order.append(shape.name)
    prs.save(out_path)
    with open(manifest_path, "w", encoding="utf-8") as fh:
        json.dump({"deck": DECK, "pages": order}, fh, indent=2, ensure_ascii=False)
        fh.write("\n")
    return order


def _add_probe_slide(prs, shape) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    box = slide.shapes.add_textbox(
        Emu(200000), Emu(200000),
        Emu(shape.width), Emu(6000000),      # original width, generous height
    )
    tf = box.text_frame
    tf.word_wrap = shape.wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left, tf.margin_right = Emu(shape.insets[0]), Emu(shape.insets[2])
    tf.margin_top, tf.margin_bottom = Emu(shape.insets[1]), Emu(shape.insets[3])

    for i, para in enumerate(shape.paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        for run in para.runs:
            if run.text == "\n":
                continue
            r = p.add_run()
            r.text = run.text
            r.font.size = Pt(run.size_pt)
            r.font.bold = run.bold
            r.font.italic = run.italic
            r.font.name = run.font


#: Where ONLYOFFICE Desktop Editors keeps its converter. The binary is the same
#: `x2t` the server products use, so a deck converted with it goes through the
#: ONLYOFFICE layout engine rather than LibreOffice's.
X2T_CANDIDATES = (
    "/Applications/ONLYOFFICE.app/Contents/Resources/converter/x2t",
    "/Applications/ONLYOFFICE.app/Contents/MacOS/converter/x2t",
    "/Applications/ONLYOFFICE Desktop Editors.app/Contents/Resources/converter/x2t",
    "/Applications/ONLYOFFICE Desktop Editors.app/Contents/MacOS/converter/x2t",
    "/opt/onlyoffice/desktopeditors/converter/x2t",
    "/usr/lib/onlyoffice/desktopeditors/converter/x2t",
    "/snap/onlyoffice-desktopeditors/current/opt/onlyoffice/desktopeditors/converter/x2t",
)


def find_x2t(hint: str = "") -> str:
    if hint and os.path.isfile(hint):
        return hint
    for c in X2T_CANDIDATES:
        if os.path.isfile(c):
            return c
    return shutil.which("x2t") or ""


def render_soffice(path: str) -> str:
    subprocess.run(
        ["soffice", "--headless", "--convert-to", "pdf",
         "--outdir", os.path.dirname(path), path],
        capture_output=True, timeout=300, check=False,
    )
    pdf = path.rsplit(".", 1)[0] + ".pdf"
    return pdf if os.path.exists(pdf) else ""


def render_x2t(path: str, binary: str) -> str:
    """Convert with ONLYOFFICE's x2t.

    Run from the converter's own directory: x2t resolves its resources - the
    DoctRenderer config, the bundled fonts - relative to where it sits, and
    fails obscurely when invoked from elsewhere. Paths are passed absolute for
    the same reason.
    """
    pdf = path.rsplit(".", 1)[0] + ".pdf"
    env = dict(os.environ)
    d = os.path.dirname(binary)
    for var in ("DYLD_LIBRARY_PATH", "DYLD_FALLBACK_LIBRARY_PATH", "LD_LIBRARY_PATH"):
        env[var] = d + (os.pathsep + env[var] if env.get(var) else "")
    r = subprocess.run(
        [binary, os.path.abspath(path), os.path.abspath(pdf)],
        capture_output=True, timeout=300, check=False, cwd=d, env=env,
    )
    if not os.path.exists(pdf):
        # x2t is quiet on success and terse on failure; surface it once rather
        # than printing "render failed" 24 times with no reason.
        render_x2t.last_error = (
            (r.stderr or b"").decode(errors="replace").strip()
            or (r.stdout or b"").decode(errors="replace").strip()
            or f"exit {r.returncode}, no output file"
        )
        return ""
    return pdf


def measure_pdf(pdf: str, tol: float = 1.5) -> tuple[int, float | None]:
    """Rendered line count, and median baseline-to-baseline pitch in points."""
    with pdfplumber.open(pdf) as doc:
        if not doc.pages:
            return 0, None
        rows: dict[float, list] = defaultdict(list)
        for ch in doc.pages[0].chars:
            rows[round(ch["top"] / tol) * tol].append(ch)
        if not rows:
            return 0, None
        tops = sorted(min(c["top"] for c in v) for v in rows.values())
    pitches = [b - a for a, b in zip(tops, tops[1:]) if b - a > 1.0]
    return len(tops), (statistics.median(pitches) if pitches else None)


def measure_pdf_pages(pdf: str, tol: float = 1.5) -> list[tuple[int, float | None]]:
    """Same measurement as measure_pdf, once per page, in page order."""
    out: list[tuple[int, float | None]] = []
    with pdfplumber.open(pdf) as doc:
        for page in doc.pages:
            rows: dict[float, list] = defaultdict(list)
            for ch in page.chars:
                rows[round(ch["top"] / tol) * tol].append(ch)
            if not rows:
                out.append((0, None))
                continue
            tops = sorted(min(c["top"] for c in v) for v in rows.values())
            pitches = [b - a for a, b in zip(tops, tops[1:]) if b - a > 1.0]
            out.append((len(tops),
                        statistics.median(pitches) if pitches else None))
    return out


def _prediction(shape):
    """The model's numbers for one shape, taken from the library, not re-derived.

    A calibration that recomputes the formula it is testing proves nothing. An
    earlier version of this file did exactly that and kept reporting a stale
    +1.7% after the model had already been fixed.
    """
    pred = layout_shape(shape)
    if pred is None:
        return None
    sizes = {r.size_pt for p in shape.paragraphs for r in p.runs if r.text.strip()}
    uniform = len(sizes) == 1
    pred_pitch = (max(sizes) * DRAWINGML_LINE_SPACING) if sizes else None
    return pred, pred_pitch, uniform


def report(label: str, deck_path: str, measured, json_path: str | None) -> int:
    """One table and one summary, whatever produced the measurements.

    Both paths - rendering shape by shape, and measuring a PDF someone exported
    by hand - end up here, so the two can never drift into reporting the same
    thing differently.
    """
    hdr = (f'{"shape":32}{"lines":13}{"pitch pred":12}{"pitch real":12}'
           f'{"pitch delta":13}')
    print(f"deck: {deck_path}   one shape per probe")
    print(f"renderer: {label}\n")
    print(hdr)
    print("-" * len(hdr))

    rows: list[dict] = []
    line_exact = line_off_by_one = line_total = 0
    pitch_deltas: list[float] = []
    untrusted: list[str] = []

    for shape, real_lines, real_pitch in measured:
        got = _prediction(shape)
        if got is None:
            continue
        pred, pred_pitch, uniform = got
        if not pred.confident:
            untrusted.append(shape.name)

        if real_lines is None:
            print(f'{shape.name[:31]:32}{"no data":13}')
            continue

        line_total += 1
        if real_lines == pred.lines:
            line_exact += 1
        elif abs(real_lines - pred.lines) == 1:
            line_off_by_one += 1

        if pred_pitch and real_pitch and uniform:
            d = (pred_pitch - real_pitch) / real_pitch
            pitch_deltas.append(abs(d))
            pitch_s, delta_s = f"{real_pitch:.2f}", f"{d * 100:+.1f}%"
        elif pred_pitch and real_pitch:
            # Mixed run sizes: each line takes its height from its own tallest
            # run, so "the pitch" is not one number and a single delta would be
            # meaningless. Line count is still comparable.
            pitch_s, delta_s = f"{real_pitch:.2f}", "mixed sizes"
        else:
            pitch_s = "-" if not real_pitch else f"{real_pitch:.2f}"
            delta_s = "single line"

        rows.append({
            "shape": shape.name,
            "slide": shape.slide,
            "lines_predicted": pred.lines,
            "lines_rendered": real_lines,
            "pitch_predicted": pred_pitch,
            "pitch_rendered": real_pitch,
            "uniform_size": uniform,
            "confident": pred.confident,
        })

        flag = "" if real_lines == pred.lines else "  <-"
        print(f'{shape.name[:31]:32}'
              f'{f"{pred.lines} / {real_lines}":13}'
              f'{f"{pred_pitch:.2f}" if pred_pitch else "-":12}'
              f'{pitch_s:12}{delta_s:13}{flag}')

    print()
    print(f"line count      exact {line_exact}/{line_total}, "
          f"off by one {line_off_by_one}/{line_total}")
    if pitch_deltas:
        pitch_deltas.sort()
        print(f"line pitch      uniform-size shapes only, n={len(pitch_deltas)}  "
              f"median {statistics.median(pitch_deltas) * 100:.2f}%  "
              f"worst {max(pitch_deltas) * 100:.2f}%")
    if untrusted:
        print(f"not confident   {len(untrusted)}: {', '.join(untrusted)}")

    if json_path:
        payload = {
            "renderer": label,
            "deck": deck_path,
            "line_count": {"exact": line_exact, "off_by_one": line_off_by_one,
                           "total": line_total},
            "pitch_delta_median": (statistics.median(pitch_deltas)
                                   if pitch_deltas else None),
            "pitch_delta_worst": max(pitch_deltas) if pitch_deltas else None,
            "shapes": rows,
        }
        with open(json_path, "w", encoding="utf-8") as fh:
            json.dump(payload, fh, indent=2, ensure_ascii=False)
            fh.write("\n")
        print(f"\nwrote {json_path}")
    return 0


def find_manifest(pdf: str) -> str:
    guess = pdf.rsplit(".", 1)[0] + ".pages.json"
    if os.path.isfile(guess):
        return guess
    hits = sorted(glob.glob(os.path.join(os.path.dirname(pdf) or ".",
                                         "*.pages.json")))
    return hits[0] if len(hits) == 1 else ""


def report_from_pdf(args) -> int:
    manifest = args.manifest or find_manifest(args.from_pdf)
    if not manifest:
        print("calibrate_pptx: cannot find the page manifest. Pass "
              "--manifest <probe>.pages.json", file=sys.stderr)
        return 2
    with open(manifest, encoding="utf-8") as fh:
        order = json.load(fh)["pages"]

    pages = measure_pdf_pages(args.from_pdf)
    if len(pages) != len(order):
        # Refuse rather than guess. A silent off-by-one here would attribute
        # every measurement to the wrong shape and the report would look fine.
        print(f"calibrate_pptx: the PDF has {len(pages)} pages but the "
              f"manifest lists {len(order)} shapes. Export the probe deck "
              f"produced by --build-probe, with every slide included.",
              file=sys.stderr)
        return 2

    by_name = {sh.name: sh for sh in read_deck(args.deck).shapes}
    missing = [n for n in order if n not in by_name]
    if missing:
        print(f"calibrate_pptx: the deck no longer has {missing[0]!r}; the "
              f"manifest was built from a different deck", file=sys.stderr)
        return 2

    measured = [(by_name[name], lines, pitch)
                for name, (lines, pitch) in zip(order, pages)]
    label = args.label or f"exported PDF ({os.path.basename(args.from_pdf)})"
    return report(label, args.deck, measured, args.json)


def main(argv: list[str] | None = None) -> int:
    ap = argparse.ArgumentParser(description=__doc__)
    ap.add_argument("--renderer", default=os.environ.get("DI_RENDERER", "soffice"),
                    help="soffice (default), x2t, or a path to a converter binary")
    ap.add_argument("--json", metavar="PATH", default=None,
                    help="also write the per-shape numbers here")
    ap.add_argument("--deck", default=DECK)
    ap.add_argument("--build-probe", metavar="PATH", default=None,
                    help="write one deck holding every probe, one shape per "
                         "slide, for renderers with no usable command line; "
                         "export it to PDF by hand, then use --from-pdf")
    ap.add_argument("--from-pdf", metavar="PATH", default=None,
                    help="measure a PDF exported from the --build-probe deck "
                         "instead of rendering; page N is the Nth shape")
    ap.add_argument("--manifest", metavar="PATH", default=None,
                    help="the <probe>.pages.json written by --build-probe")
    ap.add_argument("--label", default=None,
                    help="what produced the PDF, for the report and the JSON")
    args = ap.parse_args(argv)

    if args.build_probe:
        deck = read_deck(args.deck)
        shapes = [sh for sh in deck.shapes if layout_shape(sh) is not None]
        manifest = args.build_probe.rsplit(".", 1)[0] + ".pages.json"
        order = probe_book(shapes, args.build_probe, manifest)
        print(f"wrote {args.build_probe}  ({len(order)} slides, one per shape)")
        print(f"wrote {manifest}")
        print("\nExport that deck to PDF from the renderer you want to test, "
              "then:\n"
              "  python calibrate_pptx.py --from-pdf exported.pdf "
              "--label 'ONLYOFFICE 9.x' --json out.json")
        return 0

    if args.from_pdf:
        return report_from_pdf(args)

    name = args.renderer
    if name == "soffice":
        if not shutil.which("soffice"):
            print("calibrate_pptx: soffice is not on PATH", file=sys.stderr)
            return 2
        render = render_soffice
        label = "LibreOffice (soffice)"
    else:
        binary = find_x2t("" if name == "x2t" else name)
        if not binary:
            print("calibrate_pptx: could not find x2t. Pass its path as "
                  "--renderer /path/to/x2t", file=sys.stderr)
            return 2
        render = lambda pth: render_x2t(pth, binary)  # noqa: E731
        label = f"ONLYOFFICE x2t ({binary})"

    deck = read_deck(args.deck)
    shutil.rmtree(WORK, ignore_errors=True)
    os.makedirs(WORK, exist_ok=True)

    measured = []
    for shape in deck.shapes:
        if layout_shape(shape) is None:
            continue
        probe = os.path.join(WORK, f"{shape.name}.pptx")
        probe_deck(shape, probe)
        pdf = render(probe)
        if not pdf:
            why = getattr(render_x2t, "last_error", "")
            if why:
                print("calibrate_pptx: the converter produced nothing. Full "
                      f"message:\n{why}", file=sys.stderr)
                return 1
            measured.append((shape, None, None))
            continue
        real_lines, real_pitch = measure_pdf(pdf)
        measured.append((shape, real_lines, real_pitch))

    return report(label, args.deck, measured, args.json)


if __name__ == "__main__":
    raise SystemExit(main())
